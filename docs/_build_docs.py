"""Script para gerar Manual do Usuário (.docx) e Apresentação (.pptx)"""

# ── DOCX ──────────────────────────────────────────────────────────────────────
from docx import Document
from docx.shared import Pt, RGBColor, Inches, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
import copy

AZUL_ESCURO = RGBColor(0x1E, 0x3A, 0x5F)
AZUL_MEDIO  = RGBColor(0x2E, 0x6D, 0xA4)
AZUL_CLARO  = RGBColor(0xD6, 0xE4, 0xF0)
BRANCO      = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_ALT   = RGBColor(0xF0, 0xF4, 0xF8)


def set_cell_bg(cell, hex_color):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    shd = OxmlElement('w:shd')
    shd.set(qn('w:val'), 'clear')
    shd.set(qn('w:color'), 'auto')
    shd.set(qn('w:fill'), hex_color)
    tcPr.append(shd)


def add_heading(doc, text, level=1):
    p = doc.add_paragraph(style=f'Heading {level}')
    run = p.runs[0] if p.runs else p.add_run(text)
    if not p.runs:
        run = p.add_run(text)
    else:
        run.text = text
    run.bold = True
    run.font.color.rgb = AZUL_ESCURO
    run.font.size = Pt(14 if level == 1 else 12)
    return p


def add_body(doc, text):
    p = doc.add_paragraph()
    run = p.add_run(text)
    run.font.size = Pt(11)
    run.font.name = 'Calibri'
    return p


def add_tip(doc, text):
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.3)
    shading = OxmlElement('w:pPr')
    run = p.add_run(f'💡  {text}')
    run.font.size = Pt(10)
    run.font.color.rgb = AZUL_MEDIO
    run.font.italic = True
    return p


def add_table(doc, headers, rows, col_widths=None):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = 'Table Grid'
    table.alignment = WD_TABLE_ALIGNMENT.LEFT

    # Header row
    hdr = table.rows[0]
    for i, h in enumerate(headers):
        cell = hdr.cells[i]
        cell.text = ''
        run = cell.paragraphs[0].add_run(h)
        run.bold = True
        run.font.color.rgb = BRANCO
        run.font.size = Pt(10)
        set_cell_bg(cell, '1E3A5F')

    # Data rows
    for ri, row in enumerate(rows):
        tr = table.rows[ri + 1]
        bg = 'F0F4F8' if ri % 2 == 0 else 'FFFFFF'
        for ci, val in enumerate(row):
            cell = tr.cells[ci]
            cell.text = ''
            run = cell.paragraphs[0].add_run(str(val))
            run.font.size = Pt(10)
            set_cell_bg(cell, bg)

    # Column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            for cell in table.columns[i].cells:
                cell.width = Inches(w)

    doc.add_paragraph()
    return table


def build_docx():
    doc = Document()

    # Page margins
    for section in doc.sections:
        section.top_margin    = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin   = Inches(1.2)
        section.right_margin  = Inches(1.2)

        # Header
        header = section.header
        hp = header.paragraphs[0]
        hp.text = 'Ferragens Pinheiro — Requisições App'
        hp.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        for run in hp.runs:
            run.font.size = Pt(9)
            run.font.color.rgb = AZUL_MEDIO

        # Footer
        footer = section.footer
        fp = footer.paragraphs[0]
        fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run_f = fp.add_run('Página ')
        run_f.font.size = Pt(9)
        fldChar1 = OxmlElement('w:fldChar')
        fldChar1.set(qn('w:fldCharType'), 'begin')
        instrText = OxmlElement('w:instrText')
        instrText.text = 'PAGE'
        fldChar2 = OxmlElement('w:fldChar')
        fldChar2.set(qn('w:fldCharType'), 'end')
        run_pg = fp.add_run()
        run_pg._r.append(fldChar1)
        run_pg._r.append(instrText)
        run_pg._r.append(fldChar2)
        run_pg.font.size = Pt(9)

    # ── CAPA ──────────────────────────────────────────────────────────────────
    p_title = doc.add_paragraph()
    p_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_title.paragraph_format.space_before = Pt(80)
    run = p_title.add_run('Manual do Usuário')
    run.bold = True
    run.font.size = Pt(28)
    run.font.color.rgb = AZUL_ESCURO
    run.font.name = 'Calibri'

    p_sub = doc.add_paragraph()
    p_sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_sub.add_run('Requisições App  v1.1.0')
    run.font.size = Pt(16)
    run.font.color.rgb = AZUL_MEDIO
    run.font.name = 'Calibri'

    p_emp = doc.add_paragraph()
    p_emp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    p_emp.paragraph_format.space_before = Pt(20)
    run = p_emp.add_run('Ferragens Pinheiro')
    run.bold = True
    run.font.size = Pt(14)
    run.font.color.rgb = AZUL_ESCURO

    p_date = doc.add_paragraph()
    p_date.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = p_date.add_run('Maio de 2026')
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    doc.add_page_break()

    # ── 1. PRIMEIRO ACESSO ────────────────────────────────────────────────────
    add_heading(doc, '1. Primeiro Acesso e Login')
    add_heading(doc, 'Abrindo o sistema', 2)
    add_body(doc, 'Clique duas vezes no ícone Requisições App na área de trabalho ou na barra de tarefas.')

    add_heading(doc, 'Fazendo login', 2)
    for step in [
        '1.  Digite seu código de usuário no primeiro campo.',
        '2.  Digite sua senha no segundo campo.',
        '3.  Clique em Entrar ou pressione Enter.',
    ]:
        p = doc.add_paragraph(style='List Number')
        p.add_run(step.split('. ', 1)[1]).font.size = Pt(11)

    add_tip(doc, 'Primeiro acesso: se for a primeira vez, será solicitado que você crie uma nova senha.')

    add_heading(doc, 'Esqueci minha senha', 2)
    add_body(doc, 'Entre em contato com o administrador do sistema para redefinição de senha.')

    # ── 2. TELA PRINCIPAL ─────────────────────────────────────────────────────
    add_heading(doc, '2. Tela Principal e Navegação')
    add_heading(doc, 'Sidebar', 2)
    add_body(doc, 'A sidebar exibe apenas as telas disponíveis para o seu perfil:')
    add_table(doc,
        ['Ícone', 'Tela', 'Disponível para'],
        [
            ['📊', 'Dashboard', 'Admin, Gerente'],
            ['📋', 'Central de Pedidos', 'Todos'],
            ['🏭', 'A&R / Indústria', 'Producao, Indústria, Admin, Gerente'],
            ['👥', 'Gestão de Usuários', 'Admin'],
            ['🔧', 'Painel Técnico', 'Admin'],
            ['🕐', 'Histórico', 'Todos'],
            ['💬', 'Feedback', 'Todos'],
            ['⚙️', 'Configurações', 'Todos'],
        ],
        col_widths=[0.7, 2.0, 2.8]
    )

    add_heading(doc, 'Tema claro/escuro', 2)
    add_body(doc, 'Clique no ícone de lua/sol na parte inferior da sidebar para alternar entre modo claro e escuro.')

    # ── 3. CENTRAL DE PEDIDOS ─────────────────────────────────────────────────
    add_heading(doc, '3. Central de Pedidos')
    add_body(doc, 'A Central de Pedidos é a tela principal do sistema, onde você visualiza, cria e gerencia as requisições.')
    add_heading(doc, 'Status das requisições', 2)
    add_table(doc,
        ['Status', 'Significado'],
        [
            ['🔵 Em andamento', 'Requisição em elaboração ou aguardando'],
            ['🟡 Aguardando recebimento', 'Material aguardando chegada'],
            ['🟠 Em produção', 'Em processamento pela A&R ou Indústria'],
            ['🟢 Faturado', 'Pedido faturado e concluído'],
            ['🔴 Cancelada', 'Requisição cancelada'],
        ],
        col_widths=[2.2, 3.5]
    )

    # ── 4. CRIANDO REQUISIÇÃO ─────────────────────────────────────────────────
    add_heading(doc, '4. Criando uma Nova Requisição')
    add_tip(doc, 'Disponível para: Vendedor, Gerente, Admin.')
    add_heading(doc, 'Passo a passo', 2)
    steps = [
        'Na Central de Pedidos, clique em + Nova Requisição.',
        'Busque o cliente digitando nome, código ou CNPJ. Clique no cliente desejado.',
        'Preencha o número do pedido (campo PED).',
        'Informe se é retirada ou entrega.',
        'Adicione os itens: clique em + Adicionar Item, busque o produto, informe quantidade e preço.',
        'Adicione observações se necessário.',
        'Clique em Salvar. O PDF será gerado e enviado automaticamente para a pasta de rede.',
    ]
    for s in steps:
        p = doc.add_paragraph(style='List Number')
        p.add_run(s).font.size = Pt(11)

    # ── 5. EDITOR DE DESENHO ──────────────────────────────────────────────────
    add_heading(doc, '5. Editor de Desenho')
    add_body(doc, 'O editor permite criar croquis, anotações e medidas diretamente na requisição.')
    add_heading(doc, 'Ferramentas disponíveis', 2)
    add_table(doc,
        ['Ferramenta', 'Atalho', 'Função'],
        [
            ['Seleção', 'S', 'Seleciona e move elementos'],
            ['Caneta', 'P', 'Desenho livre'],
            ['Linha', 'L', 'Linha reta com snap nos extremos'],
            ['Seta', 'A', 'Linha com ponta de seta'],
            ['Retângulo', 'R', 'Retângulo'],
            ['Triângulo', 'T', 'Triângulo'],
            ['Texto', 'X', 'Inserir texto'],
            ['Cota MM', 'M', 'Medida em milímetros com rótulo automático'],
            ['Borracha', 'E', 'Apagar elementos'],
        ],
        col_widths=[1.5, 0.8, 3.4]
    )
    add_tip(doc, 'Pressione Esc para desmarcar a ferramenta sem fechar o editor.')
    add_tip(doc, 'A ferramenta Cota MM posiciona o rótulo automaticamente para evitar sobreposição.')

    # ── 6. HISTÓRICO ──────────────────────────────────────────────────────────
    add_heading(doc, '6. Histórico de Requisições')
    add_body(doc, 'O histórico exibe todas as requisições, independentemente do status.')
    add_heading(doc, 'Filtros disponíveis', 2)
    for item in ['Data: filtre por intervalo de datas.', 'Status: filtre por um ou mais status.', 'Vendedor/Cliente: busca por texto.']:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(item).font.size = Pt(11)

    # ── 7. NOTIFICAÇÕES ───────────────────────────────────────────────────────
    add_heading(doc, '7. Notificações')
    add_body(doc, 'O sistema envia notificações em tempo real para eventos importantes.')
    for item in [
        'O badge vermelho no sino indica notificações não lidas.',
        'Clique no sino para abrir o painel lateral com todas as notificações.',
        'Notificações novas aparecem como toasts no canto da tela.',
    ]:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(item).font.size = Pt(11)

    # ── 8. CALCULADORA ────────────────────────────────────────────────────────
    add_heading(doc, '8. Calculadora de Peso')
    add_body(doc, 'Calcula a massa de materiais metálicos com base em dimensões.')
    for step in [
        'Selecione o tipo de perfil (chato, quadrado, redondo, tubo etc.).',
        'Informe as dimensões em milímetros.',
        'Informe o comprimento em metros.',
        'O resultado em kg é exibido automaticamente.',
    ]:
        p = doc.add_paragraph(style='List Number')
        p.add_run(step).font.size = Pt(11)
    add_tip(doc, 'A densidade do aço (7,865 g/cm³) é fixa e não pode ser alterada.')

    # ── 9. CONFIGURAÇÕES ──────────────────────────────────────────────────────
    add_heading(doc, '9. Configurações')
    add_table(doc,
        ['Opção', 'Função'],
        [
            ['Tema', 'Alterna entre claro e escuro'],
            ['Tamanho da fonte', 'Ajusta o tamanho da fonte da interface'],
            ['Pasta de backgrounds', 'Define a pasta de imagens de fundo do login'],
            ['Alerta de vencimento', 'Define quantos dias antes de alertar sobre notas a vencer'],
            ['Alterar senha', 'Permite alterar sua senha de acesso'],
            ['Ver Guia Rápido', 'Abre o tour interativo do sistema'],
            ['Verificar atualizações', 'Verifica se há nova versão disponível'],
        ],
        col_widths=[2.0, 4.0]
    )

    # ── 10. GUIA RÁPIDO ───────────────────────────────────────────────────────
    add_heading(doc, '10. Guia Rápido')
    add_body(doc, 'O Guia Rápido é um tour interativo que apresenta as principais funcionalidades.')
    add_tip(doc, 'Aparece automaticamente no primeiro login. Pode ser reaberto em Configurações → Ver Guia Rápido.')
    for item in [
        'Use Próximo e Anterior para navegar entre as dicas.',
        'Clique em Fechar para encerrar o tour.',
        'O guia destaca o elemento relevante na tela com animação de foco.',
    ]:
        p = doc.add_paragraph(style='List Bullet')
        p.add_run(item).font.size = Pt(11)

    # ── 11. PERFIS ────────────────────────────────────────────────────────────
    add_heading(doc, '11. Perfis de Acesso')
    add_table(doc,
        ['Perfil', 'O que pode fazer'],
        [
            ['Admin', 'Acesso total — usuários, painel técnico, configurações avançadas'],
            ['Gerente', 'Dashboard, central de pedidos, histórico, gestão de usuários'],
            ['Vendedor', 'Criar, editar e acompanhar suas requisições'],
            ['Produção (A&R)', 'Visualizar requisições e tela de A&R em modo leitura'],
            ['Indústria', 'Visualizar requisições e tela de Indústria em modo leitura'],
            ['Entrega', 'Visualizar Central de Pedidos em modo leitura'],
        ],
        col_widths=[1.8, 4.5]
    )

    # ── 12. FAQ ───────────────────────────────────────────────────────────────
    add_heading(doc, '12. Perguntas Frequentes')
    faqs = [
        ('O PDF não foi gerado. O que fazer?',
         'Verifique se o caminho da pasta de PDFs está acessível na rede. Contate o administrador se persistir.'),
        ('Não consigo ver determinada tela.',
         'Cada perfil tem acesso apenas às telas pertinentes à sua função. Solicite acesso ao administrador.'),
        ('O sistema avisou que há uma atualização. Posso instalar?',
         'Sim. Clique em "Atualizar agora". O sistema baixa, aplica e reabre automaticamente. Seus dados são preservados.'),
        ('Esqueci minha senha.',
         'Entre em contato com o administrador do sistema para redefinição.'),
    ]
    for q, a in faqs:
        p = doc.add_paragraph()
        run = p.add_run(q)
        run.bold = True
        run.font.size = Pt(11)
        run.font.color.rgb = AZUL_ESCURO
        add_body(doc, a)
        doc.add_paragraph()

    out = r'C:\Users\João\Desktop\requisicoes\docs\Manual_do_Usuario_Requisicoes_App.docx'
    doc.save(out)
    print(f'DOCX salvo: {out}')


# ── PPTX ──────────────────────────────────────────────────────────────────────
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as PptxRGB
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

AZUL   = PptxRGB(0x1E, 0x3A, 0x5F)
AZUL_M = PptxRGB(0x2E, 0x6D, 0xA4)
AZUL_C = PptxRGB(0xD6, 0xE4, 0xF0)
BRANCO_P = PptxRGB(0xFF, 0xFF, 0xFF)
CINZA  = PptxRGB(0xF0, 0xF4, 0xF8)
VERDE  = PptxRGB(0x27, 0xAE, 0x60)
LARANJA = PptxRGB(0xE6, 0x7E, 0x22)

W = Inches(13.33)
H = Inches(7.5)


def prs_new():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def bg(slide, color: PptxRGB):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.width = line_width
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, x, y, w, h, text, size=Pt(14), color=BRANCO_P, bold=False, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = 'Calibri'
    return tb


def title_slide_layout(slide, title, subtitle='', footer='', dark=True):
    bg(slide, AZUL if dark else CINZA)
    # Title
    txbox(slide, Inches(1), Inches(2.2), Inches(11.3), Inches(1.5),
          title, size=Pt(44), color=BRANCO_P if dark else AZUL, bold=True, align=PP_ALIGN.CENTER)
    if subtitle:
        txbox(slide, Inches(1), Inches(3.8), Inches(11.3), Inches(0.8),
              subtitle, size=Pt(22), color=AZUL_C if dark else AZUL_M, align=PP_ALIGN.CENTER)
    if footer:
        txbox(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
              footer, size=Pt(14), color=PptxRGB(0xAA, 0xBB, 0xCC) if dark else AZUL_M, align=PP_ALIGN.CENTER)


def content_title(slide, title):
    txbox(slide, Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.7),
          title, size=Pt(28), color=AZUL, bold=True, align=PP_ALIGN.LEFT)
    # Divider line using a thin rectangle
    s = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(12.3), Inches(0.04))
    s.fill.solid(); s.fill.fore_color.rgb = AZUL_M
    s.line.fill.background()


def card(slide, x, y, w, h, icon, title, bullets, bg_color=AZUL_C, text_color=AZUL):
    b = box(slide, x, y, w, h, fill_color=bg_color, line_color=AZUL_M, line_width=Pt(1))
    # Icon + title
    txbox(slide, x + Inches(0.15), y + Inches(0.1), w - Inches(0.3), Inches(0.55),
          f'{icon}  {title}', size=Pt(14), color=text_color, bold=True)
    # Bullets
    tb = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.65), w - Inches(0.3), h - Inches(0.75))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b_text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f'• {b_text}'
        p.font.size = Pt(11)
        p.font.color.rgb = text_color
        p.font.name = 'Calibri'


def build_pptx():
    prs = prs_new()

    # ── SLIDE 1 — Capa ────────────────────────────────────────────────────────
    s = blank_slide(prs)
    title_slide_layout(s,
        'Requisições App',
        'Sistema de Gestão de Requisições',
        'Ferragens Pinheiro  |  v1.1.0  |  Maio de 2026'
    )

    # ── SLIDE 2 — O Desafio ───────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'O Desafio')
    challenges = [
        ('📄', 'Processos manuais', 'Requisições feitas em papel ou planilhas descentralizadas'),
        ('🔍', 'Sem rastreabilidade', 'Nenhum controle de status dos pedidos em tempo real'),
        ('⚠️', 'Erros humanos', 'PDFs gerados manualmente, sujeitos a inconsistências'),
        ('🔌', 'Equipes desconectadas', 'Produção e indústria sem acesso às informações'),
        ('📁', 'Histórico disperso', 'Consultas difíceis, dados espalhados sem padronização'),
    ]
    for i, (icon, title, desc) in enumerate(challenges):
        row = i // 3
        col = i % 3 if i < 3 else (i - 3)
        if i < 3:
            cx = Inches(0.4 + col * 4.3)
            cy = Inches(1.5)
        else:
            cx = Inches(2.55 + (i - 3) * 4.3)
            cy = Inches(4.1)
        c = box(s, cx, cy, Inches(4.0), Inches(2.2), fill_color=AZUL_C, line_color=AZUL_M, line_width=Pt(1))
        txbox(s, cx + Inches(0.15), cy + Inches(0.1), Inches(3.7), Inches(0.5),
              f'{icon}  {title}', size=Pt(13), color=AZUL, bold=True)
        txbox(s, cx + Inches(0.15), cy + Inches(0.65), Inches(3.7), Inches(1.4),
              desc, size=Pt(11), color=AZUL, wrap=True)

    # ── SLIDE 3 — A Solução ───────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'A Solução Desenvolvida')
    cols = [
        ('🖥️', 'Desktop App', ['PySide6 / Qt6', 'Interface nativa Windows', 'Escala adaptativa', 'Atualização automática']),
        ('🌐', 'API Central', ['FastAPI + Uvicorn', 'REST + SSE em tempo real', 'Autenticação JWT', 'Documentação interativa']),
        ('🗄️', 'Banco de Dados', ['PostgreSQL 12+', 'Índices GIN trigram', 'Backup automático', '12 tabelas estruturadas']),
    ]
    for i, (icon, title, bullets) in enumerate(cols):
        cx = Inches(0.5 + i * 4.27)
        card(s, cx, Inches(1.3), Inches(4.0), Inches(4.8), icon, title, bullets)
    txbox(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5),
          '100% desenvolvido internamente para as necessidades da Ferragens Pinheiro',
          size=Pt(11), color=AZUL_M, align=PP_ALIGN.CENTER)

    # ── SLIDE 4 — Perfis ─────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, '6 Perfis de Acesso')
    from pptx.util import Inches as I2
    perfis = [
        ('👑', 'Admin', 'Acesso total — usuários, painel técnico, configurações'),
        ('📊', 'Gerente', 'Dashboard, pedidos, histórico, gestão de usuários'),
        ('🛒', 'Vendedor', 'Criar e acompanhar suas requisições, histórico'),
        ('🔧', 'Produção (A&R)', 'Visualizar requisições e tela de A&R em leitura'),
        ('🏭', 'Indústria', 'Visualizar requisições e tela de Indústria em leitura'),
        ('🚚', 'Entrega', 'Visualizar Central de Pedidos em modo leitura'),
    ]
    for i, (icon, name, desc) in enumerate(perfis):
        cy = Inches(1.3 + i * 0.95)
        col = PptxRGB(0xD6, 0xE4, 0xF0) if i % 2 == 0 else BRANCO_P
        b = box(s, Inches(0.5), cy, Inches(12.3), Inches(0.88), fill_color=col)
        b.line.fill.background()
        txbox(s, Inches(0.6), cy + Inches(0.1), Inches(2.2), Inches(0.68),
              f'{icon}  {name}', size=Pt(13), color=AZUL, bold=True)
        txbox(s, Inches(2.9), cy + Inches(0.1), Inches(9.7), Inches(0.68),
              desc, size=Pt(12), color=PptxRGB(0x33, 0x33, 0x33))

    # ── SLIDE 5 — Funcionalidades ─────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'O Que o Sistema Faz')
    features = [
        ('📋', 'Central de Pedidos', 'Criação, acompanhamento e busca de requisições em tempo real'),
        ('📄', 'PDF Automático', 'Geração e envio do PDF para a pasta de rede ao salvar'),
        ('🔔', 'Notificações', 'Alertas em tempo real via SSE com toasts e painel lateral'),
        ('🖊️', 'Editor de Desenho', 'Croquis e cotas MM integrados à requisição'),
        ('📊', 'Dashboard', 'KPIs, gráficos e relatórios para gerentes e admins'),
        ('🔄', 'Auto-Update', 'Atualização automática sem intervenção de TI'),
    ]
    for i, (icon, title, desc) in enumerate(features):
        row, col_ = divmod(i, 3)
        cx = Inches(0.4 + col_ * 4.3)
        cy = Inches(1.4 + row * 2.7)
        c = box(s, cx, cy, Inches(4.1), Inches(2.4), fill_color=AZUL_C, line_color=AZUL_M, line_width=Pt(1))
        txbox(s, cx + Inches(0.15), cy + Inches(0.12), Inches(3.8), Inches(0.55),
              f'{icon}  {title}', size=Pt(14), color=AZUL, bold=True)
        txbox(s, cx + Inches(0.15), cy + Inches(0.72), Inches(3.8), Inches(1.55),
              desc, size=Pt(11), color=AZUL, wrap=True)

    # ── SLIDE 6 — Central de Pedidos ──────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Central de Pedidos')
    bullets_cp = [
        'Busca inteligente de clientes (112 mil+ registros)',
        'Status em tempo real: Em andamento → Em produção → Faturado',
        'Busca por nome, CNPJ, código ou número do pedido',
        'Ordenação por qualquer coluna com um clique',
        'PDF gerado e enviado automaticamente para a pasta de rede',
    ]
    tb = slide.shapes if False else None
    tb = s.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6.8), Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, b_text in enumerate(bullets_cp):
        if i == 0: p = tf.paragraphs[0]
        else: p = tf.add_paragraph()
        p.text = f'• {b_text}'
        p.font.size = Pt(13); p.font.color.rgb = AZUL; p.font.name = 'Calibri'
        p.space_after = Pt(8)

    # Status box (right)
    status_items = [
        (PptxRGB(0x41, 0x8B, 0xCA), '🔵 Em andamento'),
        (PptxRGB(0xF3, 0x9C, 0x12), '🟡 Aguardando recebimento'),
        (PptxRGB(0xE6, 0x7E, 0x22), '🟠 Em produção'),
        (PptxRGB(0x27, 0xAE, 0x60), '🟢 Faturado'),
        (PptxRGB(0xE7, 0x4C, 0x3C), '🔴 Cancelada'),
    ]
    hdr_b = box(s, Inches(7.7), Inches(1.3), Inches(5.1), Inches(0.55), fill_color=AZUL)
    hdr_b.line.fill.background()
    txbox(s, Inches(7.7), Inches(1.3), Inches(5.1), Inches(0.55),
          'Status das Requisições', size=Pt(13), color=BRANCO_P, bold=True, align=PP_ALIGN.CENTER)
    for i, (color, label) in enumerate(status_items):
        cy = Inches(1.85 + i * 0.82)
        b2 = box(s, Inches(7.7), cy, Inches(5.1), Inches(0.75),
                 fill_color=PptxRGB(0xF0, 0xF4, 0xF8) if i % 2 == 0 else BRANCO_P)
        b2.line.fill.background()
        dot = box(s, Inches(7.85), cy + Inches(0.2), Inches(0.35), Inches(0.35), fill_color=color)
        dot.line.fill.background()
        txbox(s, Inches(8.3), cy + Inches(0.1), Inches(4.3), Inches(0.55),
              label, size=Pt(12), color=AZUL)

    # ── SLIDE 7 — Editor de Desenho ───────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Editor de Desenho Integrado')
    items_ed = [
        ('🖊️', 'Croquis integrados', 'Desenhos e medidas diretamente na requisição'),
        ('📐', '9 Ferramentas', 'Caneta, linha, seta, retângulo, triângulo, texto, cota MM, borracha e seleção'),
        ('🎯', 'Posicionamento inteligente', 'Rótulos MM movem automaticamente para evitar sobreposição'),
        ('🔗', 'Snap de extremidades', 'Linhas se conectam precisamente a outros pontos'),
        ('📄', 'Exporta para PDF', 'Todo o desenho aparece no PDF da requisição'),
        ('⌨️', 'Atalhos de teclado', 'P, L, A, R, T, X, M, E, S — uma tecla por ferramenta'),
    ]
    for i, (icon, title, desc) in enumerate(items_ed):
        row, col_ = divmod(i, 2)
        cx = Inches(0.4 + col_ * 6.45)
        cy = Inches(1.4 + row * 1.95)
        b = box(s, cx, cy, Inches(6.1), Inches(1.75), fill_color=AZUL_C, line_color=AZUL_M, line_width=Pt(1))
        txbox(s, cx + Inches(0.15), cy + Inches(0.1), Inches(5.8), Inches(0.5),
              f'{icon}  {title}', size=Pt(13), color=AZUL, bold=True)
        txbox(s, cx + Inches(0.15), cy + Inches(0.65), Inches(5.8), Inches(0.95),
              desc, size=Pt(11), color=AZUL, wrap=True)

    # ── SLIDE 8 — Infraestrutura ──────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Infraestrutura e Confiabilidade')
    # Left column - Notificações
    b_left = box(s, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.7), fill_color=AZUL_C, line_color=AZUL_M, line_width=Pt(1))
    txbox(s, Inches(0.55), Inches(1.4), Inches(5.8), Inches(0.55),
          '🔔  Notificações em Tempo Real', size=Pt(14), color=AZUL, bold=True)
    notif_items = [
        'Tecnologia Server-Sent Events (SSE)',
        'Toasts animados + painel lateral dedicado',
        'Admins e gerentes recebem todas as notificações',
        'Reconexão automática em caso de queda de rede',
        'Badge com contador de não lidas na sidebar',
    ]
    tb2 = s.shapes.add_textbox(Inches(0.55), Inches(2.05), Inches(5.8), Inches(4.7))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    for i, item in enumerate(notif_items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f'• {item}'; p.font.size = Pt(12); p.font.color.rgb = AZUL
        p.font.name = 'Calibri'; p.space_after = Pt(8)

    # Right column - Backup
    b_right = box(s, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.7), fill_color=PptxRGB(0xE8, 0xF5, 0xE9), line_color=PptxRGB(0x27, 0xAE, 0x60), line_width=Pt(1))
    txbox(s, Inches(6.95), Inches(1.4), Inches(5.8), Inches(0.55),
          '🗄️  Backup Automático', size=Pt(14), color=PptxRGB(0x1B, 0x5E, 0x20), bold=True)
    backup_items = [
        'pg_dump agendado (horário configurável)',
        'Retenção: diário, semanal e mensal',
        'Salvo em pasta de rede protegida (TI)',
        'Configurável pelo administrador no próprio app',
        'Restauração rápida em caso de incidente',
    ]
    tb3 = s.shapes.add_textbox(Inches(6.95), Inches(2.05), Inches(5.8), Inches(4.7))
    tf3 = tb3.text_frame; tf3.word_wrap = True
    for i, item in enumerate(backup_items):
        p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
        p.text = f'• {item}'; p.font.size = Pt(12); p.font.color.rgb = PptxRGB(0x1B, 0x5E, 0x20)
        p.font.name = 'Calibri'; p.space_after = Pt(8)

    # ── SLIDE 9 — Atualização Automática ──────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Atualização Sem Intervenção Técnica')
    steps_upd = [
        ('1', 'Nova versão\npublicada\nno GitHub'),
        ('2', 'App detecta\nautomaticamente\nno startup'),
        ('3', 'Usuário clica\n"Atualizar\nagora"'),
        ('4', 'App fecha,\natualiza e\nreabre sozinho'),
    ]
    colors_upd = [AZUL, AZUL_M, PptxRGB(0x16, 0x78, 0x9A), PptxRGB(0x27, 0xAE, 0x60)]
    for i, (num, txt) in enumerate(steps_upd):
        cx = Inches(0.7 + i * 3.1)
        # Circle number
        circ = s.shapes.add_shape(9, cx + Inches(0.95), Inches(1.6), Inches(0.9), Inches(0.9))  # oval
        circ.fill.solid(); circ.fill.fore_color.rgb = colors_upd[i]; circ.line.fill.background()
        txbox(s, cx + Inches(0.95), Inches(1.6), Inches(0.9), Inches(0.9),
              num, size=Pt(18), color=BRANCO_P, bold=True, align=PP_ALIGN.CENTER)
        # Step box
        b_s = box(s, cx, Inches(2.65), Inches(2.8), Inches(2.5), fill_color=AZUL_C, line_color=colors_upd[i], line_width=Pt(2))
        txbox(s, cx + Inches(0.1), Inches(2.75), Inches(2.6), Inches(2.3),
              txt, size=Pt(13), color=AZUL, bold=True, align=PP_ALIGN.CENTER, wrap=True)
        # Arrow between steps
        if i < 3:
            txbox(s, cx + Inches(2.85), Inches(3.55), Inches(0.4), Inches(0.55),
                  '→', size=Pt(24), color=AZUL_M, align=PP_ALIGN.CENTER)

    txbox(s, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.45),
          '✅  Configurações e dados do usuário são preservados durante a atualização',
          size=Pt(12), color=PptxRGB(0x27, 0xAE, 0x60), align=PP_ALIGN.CENTER, bold=True)
    txbox(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.45),
          'Versão atual: v1.1.0  |  Versões futuras seguem o mesmo processo automático',
          size=Pt(11), color=AZUL_M, align=PP_ALIGN.CENTER)

    # ── SLIDE 10 — Resultados ─────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Resultados')
    metrics = [('6', 'Perfis de acesso'), ('12', 'Tabelas no banco'), ('11', 'Telas'), ('100%', 'Interno')]
    for i, (num, label) in enumerate(metrics):
        cx = Inches(0.4 + i * 3.2)
        m_box = box(s, cx, Inches(1.3), Inches(2.9), Inches(2.0), fill_color=AZUL, line_color=AZUL_M, line_width=Pt(1))
        txbox(s, cx, Inches(1.3), Inches(2.9), Inches(1.2),
              num, size=Pt(40), color=BRANCO_P, bold=True, align=PP_ALIGN.CENTER)
        txbox(s, cx, Inches(2.5), Inches(2.9), Inches(0.7),
              label, size=Pt(12), color=AZUL_C, align=PP_ALIGN.CENTER)
    results = [
        '✅  Rastreabilidade completa de todas as requisições',
        '✅  PDFs padronizados e organizados por vendedor automaticamente',
        '✅  Equipes de produção com acesso em tempo real',
        '✅  Backup automático do banco de dados sem intervenção',
        '✅  Atualização automática — usuário instala com um clique',
        '✅  Interface adaptável a qualquer resolução de tela',
    ]
    tb_r = s.shapes.add_textbox(Inches(0.5), Inches(3.55), Inches(12.3), Inches(3.6))
    tf_r = tb_r.text_frame; tf_r.word_wrap = True
    for i, item in enumerate(results):
        p = tf_r.paragraphs[0] if i == 0 else tf_r.add_paragraph()
        p.text = item; p.font.size = Pt(13); p.font.name = 'Calibri'
        p.font.color.rgb = PptxRGB(0x27, 0xAE, 0x60) if '✅' in item else AZUL
        p.space_after = Pt(4)

    # ── SLIDE 11 — Roadmap ────────────────────────────────────────────────────
    s = blank_slide(prs)
    bg(s, BRANCO_P)
    content_title(s, 'Próximos Passos')
    roadmap_items = [
        (PptxRGB(0x27, 0xAE, 0x60), '✅  v1.1.0 — Atual',
         'Guia Rápido por perfil • Backup automático • Editor de desenho aprimorado • Atualização automática • Notificações em tempo real'),
        (PptxRGB(0xF3, 0x9C, 0x12), '🚧  v1.2.x — Próximo',
         'Melhorias no dashboard gerencial • Integração WhatsApp • Relatórios exportáveis em Excel • Filtros avançados no histórico'),
        (PptxRGB(0x95, 0xA5, 0xA6), '🔮  Futuro',
         'App mobile para visualização • Portal web para clientes • Integração com sistema ERP • Módulo de aprovação de pedidos'),
    ]
    for i, (color, title_r, desc_r) in enumerate(roadmap_items):
        cy = Inches(1.4 + i * 1.95)
        # Left color bar
        cb = box(s, Inches(0.4), cy, Inches(0.22), Inches(1.7), fill_color=color)
        cb.line.fill.background()
        b_r = box(s, Inches(0.65), cy, Inches(12.1), Inches(1.7),
                  fill_color=PptxRGB(0xF8, 0xF9, 0xFA) if i % 2 == 0 else BRANCO_P)
        b_r.line.fill.background()
        txbox(s, Inches(0.8), cy + Inches(0.1), Inches(11.8), Inches(0.5),
              title_r, size=Pt(14), color=color, bold=True)
        txbox(s, Inches(0.8), cy + Inches(0.65), Inches(11.8), Inches(0.9),
              desc_r, size=Pt(11), color=PptxRGB(0x44, 0x44, 0x44), wrap=True)

    # ── SLIDE 12 — Encerramento ───────────────────────────────────────────────
    s = blank_slide(prs)
    title_slide_layout(s,
        'Obrigado',
        'Requisições App — desenvolvido para a Ferragens Pinheiro',
        'Dúvidas e sugestões: fale com a equipe de TI'
    )

    out = r'C:\Users\João\Desktop\requisicoes\docs\Apresentacao_Requisicoes_App.pptx'
    prs.save(out)
    print(f'PPTX salvo: {out}')


if __name__ == '__main__':
    build_docx()
    build_pptx()
    print('Concluído!')
